"""
ALTIDES (アルタイデス)
Alternative Text Insertion and Dynamic Extraction System

このプログラムは、LM Studio上のγ（gamma）モデルを利用して、PPTX、DOCX、PDFファイル内の
写真・画像、図表・グラフに対して自動的に代替テキストを生成・埋め込みします。
テキストボックスや吹き出しなど、既にテキストが含まれるオブジェクトは対象外です。

依存ライブラリ:
 - python-pptx
 - python-docx
 - PyMuPDF (fitz)
 - requests
 - Pillow
 - tkinter (標準搭載)
 - configparser (標準ライブラリ)
"""

############################
# PARAMETERS & CONFIGURATION (Configファイルから読み込み)
############################
import os
import sys
import glob
import tempfile
import logging
import configparser
import base64

# 設定ファイルの読み込み
config = configparser.ConfigParser()
config.read("config.ini")

# LLM の設定
LMSTUDIO_ENDPOINT = config.get("LLM", "endpoint", fallback="http://localhost:1234/v1/chat/completions")
LMSTUDIO_API_KEY = config.get("LLM", "api_key", fallback="")
API_TIMEOUT = config.getint("LLM", "timeout", fallback=60)
MODEL_NAME = config.get("LLM", "model", fallback="gemma-3-12b-it")

# Log file name
LOG_FILE = "altides.log"
LOG_LEVEL = config.get("Logging", "level", fallback="INFO")

# Supported file extensions
SUPPORTED_EXTENSIONS = [".pptx", ".docx", ".pdf"]

# Temporary directory for saving extracted images
TEMP_IMAGE_DIR = os.path.join(tempfile.gettempdir(), "altides_images")
os.makedirs(TEMP_IMAGE_DIR, exist_ok=True)

############################
# ライブラリのインポート
############################
import requests
from PIL import Image
from io import BytesIO

# For PPTX processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# For DOCX processing
from docx import Document

# For PDF processing
import fitz  # PyMuPDF

# For UI
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

############################
# ログ設定
############################
# 文字列からログレベル定数への変換
def get_log_level(level_name):
    levels = {
        "DEBUG": logging.DEBUG,
        "INFO": logging.INFO,
        "WARNING": logging.WARNING,
        "ERROR": logging.ERROR,
        "CRITICAL": logging.CRITICAL
    }
    return levels.get(level_name.upper(), logging.INFO)  # デフォルトはINFO

# 設定ファイルから読み取ったログレベル文字列を適切な定数に変換
log_level = get_log_level(LOG_LEVEL)

logging.basicConfig(filename=LOG_FILE, level=log_level,
                    format="%(asctime)s [%(levelname)s] %(message)s")


############################
# LM Studio γモデルを使った代替テキスト生成関数
############################
def generate_alt_text(image_bytes):
    """
    指定されたモデルを使用して画像の代替テキストを生成する
    """
    base64_image = base64.b64encode(image_bytes).decode('utf-8')
    headers = {"Content-Type": "application/json"}
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{base64_image}"}
                    },
                    {
                        "type": "text",
                        "text": "この画像の内容を英語で説明してください。Important - No talk, just go!"
                    }
                ]
            }
        ],
        "temperature": 0.7,
        "max_tokens": 4096,
        "stream": False
    }
    try:
        response = requests.post(LMSTUDIO_ENDPOINT, headers=headers, json=payload, timeout=API_TIMEOUT)
        response.raise_for_status()
        data = response.json()
        return data["choices"][0]["message"]["content"].strip()
    except Exception as e:
        logging.error(f"代替テキストの生成中にエラーが発生しました: {e}")
        return "代替テキスト生成エラー"

############################
# PPTX処理
############################
def process_pptx(file_path):
    logging.info(f"Processing PPTX: {file_path}")
    prs = Presentation(file_path)
    updated = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    # 画像のバイナリデータを取得
                    image_bytes = shape.image.blob
                    # LM Studio γモデルで代替テキストを生成
                    alt_text = generate_alt_text(image_bytes)
                    # 代替テキストを設定 (PowerPointでは「description」プロパティに設定)
                    # shape.description = alt_text # 直接設定できないため、XMLを操作
                    # 代替テキストを設定
                    shape._element.nvPicPr.cNvPr.set('descr', alt_text)
                    updated = True
                    logging.info(f"Updated PPTX image alt text: {alt_text}")
                except Exception as e:
                    logging.error(f"Error processing image in PPTX: {e}")
    if updated:
        base, ext = os.path.splitext(file_path)
        output_file = base + "_alt" + ext
        prs.save(output_file)
        logging.info(f"Saved updated PPTX: {output_file}")
        return output_file
    return None

############################
# DOCX処理
############################
def process_docx(file_path):
    logging.info(f"Processing DOCX: {file_path}")
    doc = Document(file_path)
    updated = False
    # python-docxでinline_shapeから画像情報を取得
    for shape in doc.inline_shapes:
        try:
            # 画像バイト列の取得 (内部XMLから画像の関連パートを参照)
            embed_id = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            logging.debug(f"Embed ID: {embed_id}")  # デバッグ情報：embed_id の値を確認
            image_part = doc.part.related_parts[embed_id]
            logging.debug(f"Image Part: {image_part}") # デバッグ情報：image_part オブジェクトの内容を確認
            image_bytes = image_part.blob
            alt_text = generate_alt_text(image_bytes)
            logging.debug(f"Generated Alt Text: {alt_text}")  # デバッグ情報：生成された代替テキストの値を確認
            # 代替テキストの設定：内部XMLのcNvPr要素に属性として設定
            # shape._inline.graphic.graphicData.pic.nvPicPr.cNvPr.set('descr', alt_text)
            shape._inline.docPr.set('descr', alt_text)  # 代替テキストを設定
            updated = True
            logging.info(f"Updated DOCX image alt text: {alt_text}")
        except Exception as e:
            logging.error(f"Error processing image in DOCX: {e}")
    if updated:
        base, ext = os.path.splitext(file_path)
        output_file = base + "_alt" + ext
        doc.save(output_file)
        logging.info(f"Saved updated DOCX: {output_file}")
        return output_file
    return None

############################
# PDF処理
############################
def process_pdf(file_path):
    logging.info(f"Processing PDF: {file_path}")
    doc = fitz.open(file_path)
    updated = False

    for page_num in range(len(doc)):
        page = doc[page_num]
        blocks = page.get_text("dict")["blocks"]
        logging.debug(f"Page {page_num+1}: Found {len(blocks)} blocks.")  # デバッグ情報：ページごとのブロック数を確認

        for b in blocks:
            if b["type"] == 1:  # 画像ブロック
                try:
                    width = b.get("width", 0)
                    height = b.get("height", 0)
                    xres = b.get("xres", 0)
                    yres = b.get("yres", 0)
                    logging.debug(f"Block Width: {width}, Height: {height}")  # デバッグ情報：ブロックの幅と高さを確認
                    logging.debug(f"Block XRes: {xres}, YRes: {yres}")  # デバッグ情報：ブロックのX解像度とY解像度を確認

                    # サイズと解像度の閾値を設定
                    if width < 20 or height < 20:
                        logging.debug(f"Skipping block on page {page_num+1} due to small size (width={width}, height={height}).")  # デバッグ情報：サイズが小さいためスキップされたブロックを記録
                        continue  # 小さいサイズの図形を除外
                    if xres < 36 or yres < 36:
                        logging.debug(f"Skipping block on page {page_num+1} due to low resolution (xres={xres}, yres={yres}).")  # デバッグ情報：解像度が低いためスキップされたブロックを記録
                        continue  # 低解像度の図形を除外

                    image_bytes = b["image"]
                    bbox = fitz.Rect(b["bbox"])
                    alt_text = generate_alt_text(image_bytes)
                    html_text = f"<p>{alt_text}</p>"
                    page.insert_htmlbox(bbox, html_text)
                    updated = True
                    logging.info(f"Updated PDF image alt text on page {page_num+1}: {alt_text}")
                except Exception as e:
                    logging.error(f"Error processing image in PDF on page {page_num+1}: {e}")

    if updated:
        base, ext = os.path.splitext(file_path)
        output_file = base + "_alt" + ext
        doc.save(output_file)
        logging.info(f"Saved updated PDF: {output_file}")
        return output_file

    return None


############################
# ドキュメント処理のメイン関数
############################
def process_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pptx":
        return process_pptx(file_path)
    elif ext == ".docx":
        return process_docx(file_path)
    elif ext == ".pdf":
        return process_pdf(file_path)
    else:
        logging.warning(f"Unsupported file type: {file_path}")
        return None

def process_folder(folder_path):
    processed_files = []
    for ext in SUPPORTED_EXTENSIONS:
        for file_path in glob.glob(os.path.join(folder_path, f"**/*{ext}"), recursive=True):
            output = process_file(file_path)
            if output:
                processed_files.append(output)
    return processed_files

############################
# GUI（Tkinter）: ファイル／フォルダ選択UI
############################
def run_gui():
    root = tk.Tk()
    root.title("ALTIDES - Alternative Text Insertion and Dynamic Extraction System")
    root.geometry("560x260")
    root.resizable(False, False)

    # スタイリング（ttkテーマとカスタムスタイル）
    style = ttk.Style()
    style.theme_use("clam")  # 他に "alt", "default", "classic" もある
    style.configure("TButton", padding=6, relief="flat", background="#4CAF50", foreground="white")
    style.configure("TLabel", font=("Segoe UI", 10))
    style.configure("Header.TLabel", font=("Segoe UI", 11, "bold"))
    style.configure("TEntry", padding=5)
    
    selected_path = tk.StringVar()
    output_text = tk.StringVar()

    def browse_file():
        path = filedialog.askopenfilename(
            title="Select File", filetypes=[("Supported Files", "*.pptx *.docx *.pdf")]
        )
        if path:
            selected_path.set(path)
            output_text.set("")

    def browse_folder():
        path = filedialog.askdirectory(title="Select Folder")
        if path:
            selected_path.set(path)
            output_text.set("")

    def process_action():
        path = selected_path.get()
        if not path:
            messagebox.showerror("Error", "No file or folder selected.")
            return
        processed = []
        if os.path.isdir(path):
            processed = process_folder(path)
        else:
            output = process_file(path)
            if output:
                processed.append(output)
        if processed:
            output_text.set("✅ Processing complete:\n" + "\n".join(processed))
            messagebox.showinfo("Complete", "Alternative text embedding has finished.")
        else:
            output_text.set("⚠️ The target document was not found.")
            messagebox.showwarning("Complete", "Target document not found.")

    # メインフレーム
    main_frame = ttk.Frame(root, padding="15")
    main_frame.pack(fill=tk.BOTH, expand=True)

    # ラベルと入力欄
    ttk.Label(main_frame, text="📁 Select file or folder", style="Header.TLabel").pack(anchor="w")
    entry = ttk.Entry(main_frame, textvariable=selected_path, width=65)
    entry.pack(pady=5, fill=tk.X)

    # ボタンフレーム
    button_frame = ttk.Frame(main_frame)
    button_frame.pack(pady=5, anchor="center")
    ttk.Button(button_frame, text="Select File", command=browse_file).pack(side=tk.LEFT, padx=5)
    ttk.Button(button_frame, text="Select Folder", command=browse_folder).pack(side=tk.LEFT, padx=5)

    # 実行ボタン
    ttk.Button(main_frame, text="🚀 Start Processing", command=process_action).pack(pady=10)

    # 結果表示
    ttk.Label(main_frame, textvariable=output_text, wraplength=500, justify="left").pack(pady=10, anchor="w")

    root.mainloop()

if __name__ == "__main__":
    run_gui()